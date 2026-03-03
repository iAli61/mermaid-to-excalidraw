#!/usr/bin/env python3
"""
generate_pptx.py

Parses a Marp slide deck (slides.md) and speaker notes (speaker-notes.md),
then generates a PowerPoint (.pptx) file with slides and speaker notes
using python-pptx.

Optionally exports a JSON file for the VBA macro.

Usage:
    python generate_pptx.py \
        --slides workshop/slides.md \
        --notes workshop/speaker-notes.md \
        --diagrams workshop/diagrams \
        --output workshop/presentation.pptx

    # Export JSON only (for VBA):
    python generate_pptx.py \
        --slides workshop/slides.md \
        --notes workshop/speaker-notes.md \
        --export-json workshop/slides_data.json

Dependencies:
    pip install python-pptx Pillow cairosvg
"""

import argparse
import json
import os
import re
import sys
import tempfile
from pathlib import Path

# ---------------------------------------------------------------------------
# Marp Markdown Parser
# ---------------------------------------------------------------------------

# Regex to strip Marp YAML front matter
FRONT_MATTER_RE = re.compile(r"^---\s*\n.*?\n---\s*\n", re.DOTALL)

# Regex to detect HTML comments like <!-- _class: lead -->
HTML_COMMENT_RE = re.compile(r"<!--.*?-->", re.DOTALL)

# Regex to detect image references: ![alt](path)  or  ![w:900 h:480](path)
IMAGE_RE = re.compile(r"!\[([^\]]*)\]\(([^)]+)\)")

# Regex to detect code blocks
CODE_BLOCK_RE = re.compile(r"```(\w*)\n([\s\S]*?)```")

# Regex for speaker note slide headers: ## Slide N — Title
NOTE_HEADER_RE = re.compile(r"^##\s+Slide\s+(\d+)\s*[—–-]\s*(.*)", re.MULTILINE)


def parse_slides(markdown_text: str) -> list[dict]:
    """Parse Marp markdown into a list of slide dicts."""
    # Strip front matter
    text = FRONT_MATTER_RE.sub("", markdown_text)

    # Split on slide separators (--- on its own line)
    raw_slides = re.split(r"\n---\s*\n", text)

    slides = []
    for i, raw in enumerate(raw_slides):
        raw = raw.strip()
        if not raw:
            continue

        slide = {
            "index": len(slides) + 1,
            "raw": raw,
            "title": "",
            "subtitle": "",
            "body_lines": [],
            "images": [],
            "code_blocks": [],
            "is_lead": False,
            "has_table": False,
            "speaker_notes": "",
        }

        # Detect lead class
        if "_class: lead" in raw:
            slide["is_lead"] = True
            raw = HTML_COMMENT_RE.sub("", raw).strip()

        # Strip all HTML comments
        raw = HTML_COMMENT_RE.sub("", raw).strip()

        lines = raw.split("\n")

        # Extract title (first # heading)
        body_start = 0
        for j, line in enumerate(lines):
            stripped = line.strip()
            if stripped.startswith("# ") and not stripped.startswith("## "):
                slide["title"] = stripped.lstrip("# ").strip()
                body_start = j + 1
                break
            elif stripped.startswith("## "):
                if not slide["title"]:
                    slide["title"] = stripped.lstrip("# ").strip()
                else:
                    slide["subtitle"] = stripped.lstrip("# ").strip()
                body_start = j + 1

        # Re-scan for subtitle after title
        for j in range(body_start, min(body_start + 3, len(lines))):
            stripped = lines[j].strip()
            if stripped.startswith("## "):
                slide["subtitle"] = stripped.lstrip("# ").strip()
                body_start = j + 1
                break

        # Extract images
        for match in IMAGE_RE.finditer(raw):
            slide["images"].append({
                "alt": match.group(1),
                "path": match.group(2),
            })

        # Extract code blocks
        for match in CODE_BLOCK_RE.finditer(raw):
            slide["code_blocks"].append({
                "language": match.group(1),
                "code": match.group(2).strip(),
            })

        # Extract tables as structured data
        slide["tables"] = []
        body_text = "\n".join(lines[body_start:])
        # Remove code blocks from body
        body_text_clean = CODE_BLOCK_RE.sub("", body_text)
        # Remove image references from body
        body_text_clean = IMAGE_RE.sub("", body_text_clean)

        # Parse markdown tables
        body_no_tables = []
        table_lines_buf = []
        in_table = False
        for ln in body_text_clean.split("\n"):
            stripped = ln.strip()
            if stripped.startswith("|") and stripped.endswith("|"):
                in_table = True
                table_lines_buf.append(stripped)
            else:
                if in_table:
                    # End of a table block — parse it
                    parsed = _parse_markdown_table(table_lines_buf)
                    if parsed:
                        slide["tables"].append(parsed)
                    table_lines_buf = []
                    in_table = False
                body_no_tables.append(ln)
        # Flush any trailing table
        if table_lines_buf:
            parsed = _parse_markdown_table(table_lines_buf)
            if parsed:
                slide["tables"].append(parsed)

        slide["has_table"] = len(slide["tables"]) > 0

        # Clean up body lines (tables already removed)
        slide["body_lines"] = [
            ln for ln in body_no_tables
            if ln.strip() and not ln.strip().startswith("## ")
        ]

        slides.append(slide)

    return slides


def _parse_markdown_table(lines: list[str]) -> dict | None:
    """Parse a list of markdown table lines into headers + rows."""
    if len(lines) < 2:
        return None
    # Split cells, strip whitespace and outer pipes
    def split_row(line):
        cells = line.strip().strip("|").split("|")
        return [c.strip() for c in cells]

    headers = split_row(lines[0])
    # Skip separator row(s) that contain only dashes/colons/pipes
    data_rows = []
    for ln in lines[1:]:
        cells = split_row(ln)
        # Separator rows look like: ----, :---, :---:, etc.
        if all(re.match(r'^[:\-\s]+$', c) for c in cells):
            continue
        data_rows.append(cells)
    if not headers:
        return None
    return {"headers": headers, "rows": data_rows}


def parse_speaker_notes(notes_text: str) -> dict[int, str]:
    """Parse speaker notes markdown into a dict of slide_index -> notes_text."""
    notes = {}
    # Split on ## Slide N headers
    sections = NOTE_HEADER_RE.split(notes_text)

    # sections = [preamble, index1, title1, content1, index2, title2, content2, ...]
    i = 1
    while i < len(sections) - 2:
        slide_num = int(sections[i])
        # title = sections[i + 1]
        content = sections[i + 2].strip()

        # Clean up content: remove --- separators, keep plain text
        content = re.sub(r"\n---\s*$", "", content).strip()
        notes[slide_num] = content
        i += 3

    return notes


# ---------------------------------------------------------------------------
# PPTX Generation
# ---------------------------------------------------------------------------

def create_pptx(slides: list[dict], diagrams_dir: str, output_path: str):
    """Generate a .pptx file from parsed slides."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        from pptx.dml.color import RGBColor
    except ImportError:
        print("ERROR: python-pptx is required. Install with: pip install python-pptx")
        sys.exit(1)

    prs = Presentation()
    # Set widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme matching Marp theme
    TITLE_COLOR = RGBColor(0x2E, 0x50, 0x90)
    SUBTITLE_COLOR = RGBColor(0x55, 0x55, 0x55)
    BODY_COLOR = RGBColor(0x33, 0x33, 0x33)
    ACCENT_COLOR = RGBColor(0xD9, 0x53, 0x4F)
    CODE_BG = RGBColor(0xF5, 0xF5, 0xF5)

    for slide_data in slides:
        if slide_data["is_lead"] and slide_data["title"]:
            # Title slide layout
            layout = prs.slide_layouts[0]  # Title Slide
            slide = prs.slides.add_slide(layout)

            # Title
            if slide.shapes.title:
                slide.shapes.title.text = slide_data["title"]
                for para in slide.shapes.title.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = TITLE_COLOR
                        run.font.size = Pt(44)
                        run.font.bold = True

            # Subtitle
            if len(slide.placeholders) > 1:
                subtitle_ph = slide.placeholders[1]
                body_text = "\n".join(
                    ln.strip().lstrip("*").rstrip("*").strip()
                    for ln in slide_data["body_lines"]
                    if ln.strip()
                )
                subtitle_ph.text = body_text or slide_data.get("subtitle", "")
                for para in subtitle_ph.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = SUBTITLE_COLOR
                        run.font.size = Pt(24)

        else:
            # Content slide layout
            layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(layout)

            # Title
            if slide.shapes.title:
                title_text = slide_data["title"]
                if slide_data.get("subtitle"):
                    title_text += " — " + slide_data["subtitle"]
                slide.shapes.title.text = title_text
                for para in slide.shapes.title.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = TITLE_COLOR
                        run.font.size = Pt(32)
                        run.font.bold = True

            # Body content
            if len(slide.placeholders) > 1:
                body_ph = slide.placeholders[1]
                tf = body_ph.text_frame
                tf.clear()

                # Add body lines (non-table text)
                first = True
                for line in slide_data["body_lines"]:
                    clean = line.strip()
                    if not clean:
                        continue

                    # Detect bullet level
                    level = 0
                    if clean.startswith("- "):
                        clean = clean[2:]
                    elif clean.startswith("  - "):
                        clean = clean[4:]
                        level = 1

                    # Strip markdown bold
                    is_bold = "**" in clean
                    clean = clean.replace("**", "")

                    # Strip markdown emphasis (our custom red)
                    has_emphasis = clean.startswith("*") and clean.endswith("*")
                    clean = clean.strip("*")

                    # Strip blockquote
                    if clean.startswith("> "):
                        clean = clean[2:]

                    if first:
                        para = tf.paragraphs[0]
                        first = False
                    else:
                        para = tf.add_paragraph()

                    para.text = clean
                    para.level = level
                    para.alignment = PP_ALIGN.LEFT

                    for run in para.runs:
                        run.font.size = Pt(20)
                        run.font.color.rgb = BODY_COLOR
                        if is_bold:
                            run.font.bold = True
                        if has_emphasis:
                            run.font.color.rgb = ACCENT_COLOR

                # Add code blocks as text
                for code_block in slide_data["code_blocks"]:
                    para = tf.add_paragraph()
                    para.text = ""  # spacer
                    code_para = tf.add_paragraph()
                    code_para.text = code_block["code"][:500]  # Truncate long code
                    for run in code_para.runs:
                        run.font.name = "Consolas"
                        run.font.size = Pt(14)
                        run.font.color.rgb = BODY_COLOR

            # Render tables as native PPTX table objects
            for tbl_data in slide_data.get("tables", []):
                _add_pptx_table(
                    slide, tbl_data, Inches, Pt, RGBColor,
                    TITLE_COLOR, BODY_COLOR,
                )

            # Add images
            for img_info in slide_data["images"]:
                img_path = os.path.join(diagrams_dir, os.path.basename(img_info["path"]))
                if not os.path.exists(img_path):
                    # Try relative to diagrams_dir parent
                    img_path = os.path.join(
                        os.path.dirname(diagrams_dir),
                        img_info["path"]
                    )

                if os.path.exists(img_path):
                    # Convert SVG to PNG if needed
                    final_path = img_path
                    if img_path.lower().endswith(".svg"):
                        final_path = _svg_to_png(img_path)
                        if not final_path:
                            continue

                    try:
                        # Center the image
                        slide.shapes.add_picture(
                            final_path,
                            Inches(1.5),
                            Inches(2.0),
                            width=Inches(10),
                        )
                    except Exception as e:
                        print(f"  ⚠ Could not add image {img_path}: {e}")

        # Add speaker notes
        if slide_data.get("speaker_notes"):
            notes_slide = slide.notes_slide
            notes_tf = notes_slide.notes_text_frame
            notes_tf.text = slide_data["speaker_notes"]

    prs.save(output_path)
    print(f"✓ PPTX saved: {output_path} ({len(slides)} slides)")


def _add_pptx_table(slide, tbl_data, Inches, Pt, RGBColor,
                    header_color, body_color):
    """Add a native PPTX table to the slide from parsed table data."""
    from pptx.util import Emu
    from pptx.enum.text import PP_ALIGN

    headers = tbl_data["headers"]
    rows = tbl_data["rows"]
    n_cols = len(headers)
    n_rows = 1 + len(rows)  # header + data rows

    if n_cols == 0 or n_rows == 0:
        return

    # Position: centered horizontally, below the title area
    # Determine available body space
    slide_w = slide.part.package.presentation_part.presentation.slide_width
    tbl_width = int(slide_w * 0.85)
    left = int((slide_w - tbl_width) / 2)
    top = Inches(2.0)
    row_height = Inches(0.45)
    tbl_height = int(row_height * n_rows)

    # Remove the body placeholder if it's empty (table replaces it)
    # Keep it if it has real content
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == 1:  # body placeholder
            tf = ph.text_frame
            text = tf.text.strip()
            if not text:
                sp = ph._element
                sp.getparent().remove(sp)
            break

    table_shape = slide.shapes.add_table(
        n_rows, n_cols, left, top, tbl_width, tbl_height
    )
    table = table_shape.table

    # Style constants
    HEADER_BG = RGBColor(0x2E, 0x50, 0x90)
    HEADER_FG = RGBColor(0xFF, 0xFF, 0xFF)
    EVEN_BG = RGBColor(0xF2, 0xF2, 0xF2)
    ODD_BG = RGBColor(0xFF, 0xFF, 0xFF)
    BORDER_COLOR = RGBColor(0xBF, 0xBF, 0xBF)
    FONT_SIZE_HEADER = Pt(14)
    FONT_SIZE_BODY = Pt(13)

    # Distribute column widths evenly
    col_width = int(tbl_width / n_cols)
    for ci in range(n_cols):
        table.columns[ci].width = col_width

    # Fill header row
    for ci, hdr in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = hdr.replace("**", "").replace("`", "")
        for para in cell.text_frame.paragraphs:
            para.alignment = PP_ALIGN.LEFT
            for run in para.runs:
                run.font.size = FONT_SIZE_HEADER
                run.font.bold = True
                run.font.color.rgb = HEADER_FG
                run.font.name = "Calibri"
        # Header background
        _set_cell_bg(cell, HEADER_BG)

    # Fill data rows
    for ri, row_cells in enumerate(rows):
        bg = EVEN_BG if ri % 2 == 0 else ODD_BG
        for ci in range(n_cols):
            cell = table.cell(ri + 1, ci)
            value = row_cells[ci] if ci < len(row_cells) else ""
            # Strip markdown formatting
            value = value.replace("**", "").replace("`", "")
            cell.text = value
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.size = FONT_SIZE_BODY
                    run.font.color.rgb = body_color
                    run.font.name = "Calibri"
            _set_cell_bg(cell, bg)


def _set_cell_bg(cell, rgb_color):
    """Set the background fill of a PPTX table cell."""
    from pptx.oxml.ns import qn
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    srgbClr = solidFill.makeelement(qn('a:srgbClr'), {
        'val': '%02X%02X%02X' % (rgb_color[0], rgb_color[1], rgb_color[2])
    })
    solidFill.append(srgbClr)
    tcPr.append(solidFill)


def _svg_to_png(svg_path: str) -> str | None:
    """Convert SVG to PNG using cairosvg. Returns PNG path or None."""
    try:
        import cairosvg
        png_path = svg_path.rsplit(".", 1)[0] + ".png"
        if not os.path.exists(png_path):
            cairosvg.svg2png(
                url=svg_path,
                write_to=png_path,
                output_width=1920,
            )
        return png_path
    except ImportError:
        print(f"  ⚠ cairosvg not installed — cannot convert {svg_path}")
        print("    Install with: pip install cairosvg")
        return None
    except Exception as e:
        print(f"  ⚠ SVG→PNG conversion failed for {svg_path}: {e}")
        return None


# ---------------------------------------------------------------------------
# JSON Export (for VBA)
# ---------------------------------------------------------------------------

def export_json(slides: list[dict], output_path: str):
    """Export slides data as JSON for the VBA macro."""
    json_slides = []
    for s in slides:
        json_slides.append({
            "index": s["index"],
            "title": s["title"],
            "subtitle": s.get("subtitle", ""),
            "is_lead": s["is_lead"],
            "body_lines": [
                ln.strip().replace("**", "").lstrip("- ").lstrip("> ")
                for ln in s["body_lines"] if ln.strip()
            ],
            "images": s["images"],
            "code_blocks": [cb["code"] for cb in s["code_blocks"]],
            "has_table": s["has_table"],
            "table_raw": _extract_table(s["raw"]) if s["has_table"] else "",
            "speaker_notes": s.get("speaker_notes", ""),
        })

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump({"slides": json_slides}, f, indent=2, ensure_ascii=False)

    print(f"✓ JSON exported: {output_path} ({len(json_slides)} slides)")


def _extract_table(raw: str) -> str:
    """Extract the first Markdown table from a raw slide."""
    lines = raw.split("\n")
    table_lines = []
    in_table = False
    for ln in lines:
        if "|" in ln and ln.strip().startswith("|"):
            in_table = True
            table_lines.append(ln.strip())
        elif in_table:
            break
    return "\n".join(table_lines)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    parser = argparse.ArgumentParser(
        description="Generate PPTX from Marp slides + speaker notes"
    )
    parser.add_argument("--slides", required=True, help="Path to Marp slides.md")
    parser.add_argument("--notes", help="Path to speaker-notes.md")
    parser.add_argument("--diagrams", default="", help="Path to diagrams directory")
    parser.add_argument("--output", default="presentation.pptx", help="Output .pptx path")
    parser.add_argument("--export-json", dest="json_path", help="Export JSON for VBA macro")

    args = parser.parse_args()

    # Read slides
    slides_text = Path(args.slides).read_text(encoding="utf-8")
    slides = parse_slides(slides_text)
    print(f"Parsed {len(slides)} slides from {args.slides}")

    # Read speaker notes
    if args.notes and os.path.exists(args.notes):
        notes_text = Path(args.notes).read_text(encoding="utf-8")
        notes_map = parse_speaker_notes(notes_text)
        print(f"Parsed speaker notes for {len(notes_map)} slides from {args.notes}")

        # Merge notes into slides
        for slide in slides:
            idx = slide["index"]
            if idx in notes_map:
                slide["speaker_notes"] = notes_map[idx]
    else:
        print("No speaker notes provided (or file not found)")

    # Resolve diagrams directory
    diagrams_dir = args.diagrams
    if not diagrams_dir:
        diagrams_dir = os.path.join(os.path.dirname(args.slides), "diagrams")

    # Export JSON
    if args.json_path:
        export_json(slides, args.json_path)

    # Generate PPTX
    if not args.json_path or args.output != "presentation.pptx":
        create_pptx(slides, diagrams_dir, args.output)


if __name__ == "__main__":
    main()
